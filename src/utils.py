from pptx import Presentation

def remove_all_slides(prs: Presentation):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    for slide in slides:
        xml_slides.remove(slide)
    print("所有默认幻灯片已被移除。")

def remove_slides(prs: Presentation, count=0):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    rm_count = max(0, min(count, len(slides)))
    for i in range(0, rm_count):
        slide = slides[i]
        xml_slides.remove(slide)
    print(f"所有默认{rm_count}张幻灯片已被移除。")

def get_slides_count(prs: Presentation):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    return len(slides)