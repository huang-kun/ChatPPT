from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
import os

'''
# ChatPPT Demo

## 2024 业绩概述
- 总收入增长15%
- 市场份额扩大至30%

## 业绩图表
- OpenAI 利润不断增加
![业绩图表](images/performance_chart.png)

## 新产品发布
- 产品A: 特色功能介绍
- 产品B: 市场定位
![未来增长](images/forecast.png)
'''

'''
0: Title Slide
1: Title and Content
2: Section Header
3: Two Content
4: Comparison
5: Title Only
6: Blank
7: Content with Caption
8: Picture with Caption
    0: TITLE (1)
    1: PICTURE (18)
    2: BODY (2)
    10: DATE (16)
    11: FOOTER (15)
    12: SLIDE_NUMBER (13)
9: Title and Vertical Text
10: Vertical Title and Text
'''


def get_path(rel_path):
    curr_path = os.path.abspath(__file__)
    src_dir = os.path.dirname(curr_path)
    project_dir = os.path.dirname(src_dir)
    return os.path.join(project_dir, rel_path)


def create_homework_ppt():
    prs = Presentation()
    layout_mapping = {}

    # Get Layouts
    for i, layout in enumerate(prs.slide_layouts):
        layout_mapping[layout.name] = i

    create_title_slide(prs, layout_mapping, "ChatPPT Demo", "python-pptx homework")
    create_content_slide(prs, layout_mapping, "2024 业绩概述", ["总收入增长15%", "市场份额扩大至30%"])
    create_image_slide(prs, layout_mapping, "业绩图表", ["OpenAI 利润不断增加"], "performance_chart.png")
    create_image_slide(prs, layout_mapping, "新产品发布", ["产品A: 特色功能介绍", "产品B: 市场定位"], "forecast.png")

    # Save
    outputs_dir = get_path('outputs')
    if not os.path.exists(outputs_dir):
        os.mkdir(outputs_dir)
    
    prs.save(get_path('outputs' + os.sep + 'ChatPPT_Manual_Demo.pptx'))


def create_title_slide(prs, layout_mapping, title, subtitle):
    idx = layout_mapping.get("Title Slide", 0)
    title_slide_layout = prs.slide_layouts[idx]
    title_slide = prs.slides.add_slide(title_slide_layout)

    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]

    title.text = "ChatPPT Demo"
    subtitle.text = "python-pptx homework"


def create_content_slide(prs, layout_mapping, title, bullets):
    idx = layout_mapping.get("Title and Content", 1)
    content_layout = prs.slide_layouts[idx]
    content_slide = prs.slides.add_slide(content_layout)
    
    shapes = content_slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = title
    
    for bullet in bullets:
        text_frame = body_shape.text_frame
        p = text_frame.add_paragraph()
        p.text = bullet
        p.level = 0


def create_image_slide(prs, layout_mapping, title, bullets, image_name):
    idx = layout_mapping.get("Picture with Caption", 8)
    image_layout = prs.slide_layouts[idx]
    image_slide = prs.slides.add_slide(image_layout)

    shapes = image_slide.shapes
    for shape in shapes:
        if shape.is_placeholder:
            phf = shape.placeholder_format
            if phf.type == PP_PLACEHOLDER.TITLE:
                shape.text_frame.text = title
            elif phf.type == PP_PLACEHOLDER.BODY:
                for bullet in bullets:
                    p = shape.text_frame.add_paragraph()
                    p.text = bullet
                    p.level = 1
            elif phf.type == PP_PLACEHOLDER.PICTURE:
                image_path = get_path("images" + os.sep + image_name)
                shape.insert_picture(image_path)


if __name__ == "__main__":
    create_homework_ppt()

